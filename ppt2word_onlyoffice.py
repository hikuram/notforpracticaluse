# ==============================================================================
# PowerPoint to Word "Total Victory V4" Converter 
# (Multi-file, Headers, Hidden Slide Removal, Absolute Table Width & Height)
# ==============================================================================
import os
import sys
import subprocess
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.oxml import parse_xml, OxmlElement
from docx.oxml.ns import qn
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ROW_HEIGHT_RULE
from PIL import Image
from pdf2image import convert_from_path

def export_slides_to_images_onlyoffice(ppt_path, output_dir):
    ppt_abs_path = os.path.abspath(ppt_path)
    base_name = os.path.splitext(os.path.basename(ppt_path))[0]
    pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
    
    builder_script_path = os.path.join(output_dir, "convert.docbuilder")
    with open(builder_script_path, "w", encoding="utf-8") as f:
        f.write(f'builder.OpenFile("{ppt_abs_path}");\n')
        f.write(f'builder.SaveFile("pdf", "{pdf_path}");\n')
        f.write('builder.CloseFile();\n')

    cmd = ["documentbuilder", builder_script_path]
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except subprocess.CalledProcessError as e:
        print(f"ONLYOFFICE変換エラー ({base_name}): {e}")
        sys.exit(1)
    
    images = convert_from_path(pdf_path, dpi=300)
    image_paths = []
    for i, image in enumerate(images):
        img_path = os.path.join(output_dir, f"{base_name}_slide_{i:03d}.jpg")
        # RGBに変換した上でJPEGとして保存（quality=85がおすすめ）
        image.convert('RGB').save(img_path, "JPEG", quality=85)
        image_paths.append(img_path)
        
    return image_paths

def pad_image_to_16_9(image_path):
    img = Image.open(image_path)
    width, height = img.size
    target_ratio = 16.0 / 9.0
    if abs(width / height - target_ratio) > 0.05:
        if width / height < target_ratio:
            new_width = int(height * target_ratio)
            # RGBに変換した上でJPEGとして保存（quality=85がおすすめ）
            new_img = Image.new("RGB", (new_width, height), (255, 255, 255))
            new_img.paste(img, ((new_width - width) // 2, 0))
            new_img.save(image_path, "JPEG", quality=85)

def is_slide_visible(slide):
    show = slide.element.get('show')
    return show != '0' and show != 'false'

def extract_text_from_slide(slide):
    text_runs = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        text_runs.append(run.text.strip())
    return " ".join(text_runs)

def apply_collapse_hack(paragraph):
    try:
        pPr = paragraph._p.get_or_add_pPr()
        xml_str = '<w15:collapsed w:val="1" xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:w15="http://schemas.microsoft.com/office/word/2012/wordml"/>'
        pPr.append(parse_xml(xml_str))
    except Exception:
        pass

def set_style_font_to_meiryo(doc):
    target_styles = ['Normal', 'Title', 'Heading 1', 'Heading 2', 'Heading 3', 'Table Grid']
    for style_name in target_styles:
        try:
            style = doc.styles[style_name]
            style.font.name = 'メイリオ'
            rPr = style.element.get_or_add_rPr()
            rFonts = rPr.xpath('w:rFonts')
            if rFonts:
                rFonts[0].set(qn('w:eastAsia'), 'メイリオ')
            else:
                rFonts = OxmlElement('w:rFonts')
                rFonts.set(qn('w:eastAsia'), 'メイリオ')
                rPr.append(rFonts)
        except KeyError:
            pass

def set_table_border_color(table, color="D3D3D3"):
    """表の枠線を指定した色（HEXコード）に変更するハック"""
    tblPr = table._tbl.tblPr
    tblBorders = tblPr.find(qn('w:tblBorders'))
    if tblBorders is not None:
        tblPr.remove(tblBorders)
        
    new_tblBorders = OxmlElement('w:tblBorders')
    for b in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        border = OxmlElement(f'w:{b}')
        border.set(qn('w:val'), 'single')
        border.set(qn('w:sz'), '4')  # 線の太さ（4 = 0.5pt）
        border.set(qn('w:space'), '0')
        border.set(qn('w:color'), color)
        new_tblBorders.append(border)
    tblPr.append(new_tblBorders)

def enforce_absolute_table_width(table, left_cm, right_cm):
    table.autofit = False
    tblPr = table._tbl.tblPr
    tblLayout = OxmlElement('w:tblLayout')
    tblLayout.set(qn('w:type'), 'fixed')
    tblPr.append(tblLayout)
    
    def set_cell_width(cell, width_cm):
        tcPr = cell._tc.get_or_add_tcPr()
        tcW = OxmlElement('w:tcW')
        tcW.set(qn('w:w'), str(int(width_cm * 567)))
        tcW.set(qn('w:type'), 'dxa')
        tcPr.append(tcW)

    set_cell_width(table.rows[0].cells[0], left_cm)
    set_cell_width(table.rows[0].cells[1], right_cm)

def process_multiple_files(ppt_paths, output_word_path):
    temp_dir = os.path.join(os.path.dirname(ppt_paths[0]), "_ppt_temp_v4")
    os.makedirs(temp_dir, exist_ok=True)
    
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(2.0)
        section.right_margin = Cm(2.0)
    set_style_font_to_meiryo(doc)

    # --- 全体ヘッダー (Title) ---
    p_main = doc.add_paragraph("■ 統合議事録")
    p_main.style = 'Title'
    p_main.runs[0].font.size = Pt(16)
    
    table_meta = doc.add_table(rows=2, cols=2)
    table_meta.style = 'Table Grid'
    enforce_absolute_table_width(table_meta, 3.0, 14.0)
    
    # 16:9画像の高さ(約6.2cm)に合わせ、2行で合計6.2cmの高さを確保
    for row in table_meta.rows:
        row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
        row.height = Cm(3.0) 
        
    table_meta.cell(0, 0).text = "日時："
    table_meta.cell(1, 0).text = "参加者："
    doc.add_paragraph().paragraph_format.space_after = Pt(24)

    # --- 各ファイルの処理ループ ---
    for ppt_path in ppt_paths:
        base_name = os.path.basename(ppt_path)
        print(f"\n処理中: {base_name}")
        
        image_files = export_slides_to_images_onlyoffice(ppt_path, temp_dir)
        for img_file in image_files:
            pad_image_to_16_9(img_file)
            
        prs = Presentation(ppt_path)
        visible_slides = [s for s in prs.slides if is_slide_visible(s)]
        
        if len(image_files) == len(visible_slides):
            target_slides = visible_slides
            target_images = image_files
        elif len(image_files) == len(prs.slides):
            target_slides = visible_slides
            target_images = [image_files[i] for i, s in enumerate(prs.slides) if is_slide_visible(s)]
        else:
            print(f"警告: {base_name} のスライド数と画像数が一致しません。全て出力します。")
            target_slides = list(prs.slides)
            target_images = image_files

        # --- ファイルごとヘッダーブロック (Heading 1) ---
        p_file = doc.add_paragraph(f"▼ 資料: {base_name}")
        p_file.style = 'Heading 1'
        p_file.runs[0].font.color.rgb = RGBColor(0, 51, 153)
        p_file.paragraph_format.keep_with_next = True
        
        table_file = doc.add_table(rows=1, cols=2)
        table_file.style = 'Table Grid'
        enforce_absolute_table_width(table_file, 3.0, 14.0)
        
        # スライド画像と同じ高さ(6.2cm)を確保
        for row in table_file.rows:
            row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
            row.height = Cm(6.0)
            
        table_file.cell(0, 0).text = "特記事項："
        
        for row in table_file.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    paragraph.paragraph_format.keep_with_next = True
                    
        doc.add_paragraph().paragraph_format.space_after = Pt(12)

        # --- スライドごとの出力 (Heading 2) ---
        for i, (slide, img_path) in enumerate(zip(target_slides, target_images)):
            slide_text = extract_text_from_slide(slide)
            
            p_title = doc.add_paragraph()
            p_title.style = 'Heading 2'
            p_title.paragraph_format.keep_with_next = True
            run_title = p_title.add_run(f"■ {base_name} - P.{i+1}")
            run_title.font.size = Pt(10)
            run_title.font.color.rgb = RGBColor(255, 255, 255)

            table = doc.add_table(rows=1, cols=2)
            table.style = 'Table Grid'
            enforce_absolute_table_width(table, 11.5, 5.5)
            set_table_border_color(table, "B0B0B0")

            cell_left = table.rows[0].cells[0]
            cell_right = table.rows[0].cells[1]

            p_img = cell_left.paragraphs[0]
            p_img.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p_img.add_run().add_picture(img_path, width=Cm(11.0))
            
            p_note = cell_right.paragraphs[0]
            
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        paragraph.paragraph_format.keep_with_next = True

            # --- 抽出テキスト (Heading 3) ---
            if slide_text:
                p_summary = doc.add_paragraph()
                p_summary.style = 'Heading 3'
                p_summary.paragraph_format.keep_with_next = True
                run_summary = p_summary.add_run("▶ 抽出テキスト（検索・コピペ用）")
                run_summary.font.size = Pt(9)
                run_summary.font.color.rgb = RGBColor(255, 255, 255)
                apply_collapse_hack(p_summary)
                
                p_text = doc.add_paragraph(slide_text)
                p_text.style = 'Normal'
                for run in p_text.runs:
                    run.font.size = Pt(8.5)
                    run.font.color.rgb = RGBColor(80, 80, 80)
                    
            p_spacer = doc.add_paragraph()
            p_spacer.paragraph_format.space_after = Pt(12)

    print(f"\n全ファイルの処理が完了しました。保存中: {os.path.basename(output_word_path)}")
    doc.save(output_word_path)
    
    import shutil
    shutil.rmtree(temp_dir)
    print("完了！統合議事録が生成されました。")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("使い方: python ppt2word_v4.py <PPTXファイル1> [PPTXファイル2 ...]")
        sys.exit(1)
        
    ppt_files = sys.argv[1:]
    base_dir = os.path.dirname(os.path.abspath(ppt_files[0]))
    out_name = "統合議事録ベース.docx"
    output_word = os.path.join(base_dir, out_name)
    
    process_multiple_files(ppt_files, output_word)

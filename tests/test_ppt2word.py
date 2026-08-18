import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
import ppt2word


class XmlTextSanitizationTests(unittest.TestCase):
    def test_replaces_soft_breaks_and_removes_invalid_xml_characters(self):
        text = "A\x00B\vC\x08D\tE\rF\nG"
        self.assertEqual("AB\nCD\tE\rF\nG", ppt2word.sanitize_xml_text(text))


class TextExtractionTests(unittest.TestCase):
    def test_extracts_soft_break_as_newline(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4.0), Inches(1.0)
        )
        text_box.text_frame.text = "First\vSecond"

        extracted = ppt2word.extract_text_from_slide(slide)
        self.assertEqual("First\nSecond", extracted)

    def test_extracts_runs_tables_and_group_shapes(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4.0), Inches(1.0)
        )
        paragraph = text_box.text_frame.paragraphs[0]
        paragraph.add_run().text = "Power"
        paragraph.add_run().text = "Point"

        table_shape = slide.shapes.add_table(
            1, 1, Inches(0.5), Inches(2.0), Inches(3.0), Inches(1.0)
        )
        table_shape.table.cell(0, 0).text = "TableText"

        group_shape = slide.shapes.add_group_shape()
        group_text_box = group_shape.shapes.add_textbox(
            Inches(0.5), Inches(3.5), Inches(3.0), Inches(1.0)
        )
        group_text_box.text_frame.text = "GroupedText"

        extracted = ppt2word.extract_text_from_slide(slide)
        self.assertIn("PowerPoint", extracted)
        self.assertNotIn("Power Point", extracted)
        self.assertIn("TableText", extracted)
        self.assertIn("GroupedText", extracted)


class AspectRatioTests(unittest.TestCase):
    def test_contains_narrow_image_without_cropping(self):
        with tempfile.TemporaryDirectory() as directory:
            image_path = Path(directory) / "narrow.jpg"
            Image.new("RGB", (1200, 900), "white").save(image_path, "JPEG")

            ppt2word.pad_image_to_16_9(image_path)

            with Image.open(image_path) as image:
                self.assertEqual((1600, 900), image.size)

    def test_contains_wide_image_without_cropping(self):
        with tempfile.TemporaryDirectory() as directory:
            image_path = Path(directory) / "wide.jpg"
            Image.new("RGB", (2100, 900), "white").save(image_path, "JPEG")

            ppt2word.pad_image_to_16_9(image_path)

            with Image.open(image_path) as image:
                self.assertEqual(2100, image.width)
                self.assertGreater(image.height, 900)
                self.assertAlmostEqual(16.0 / 9.0, image.width / image.height, places=3)


class PdfMappingTests(unittest.TestCase):
    def make_presentation_with_hidden_slide(self):
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        hidden_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        hidden_slide.element.set("show", "0")
        presentation.slides.add_slide(presentation.slide_layouts[6])
        return presentation

    def test_maps_pdf_that_contains_all_slides(self):
        presentation = self.make_presentation_with_hidden_slide()
        pairs = ppt2word.resolve_slide_page_pairs(presentation, 3, "sample.pptx")
        self.assertEqual([(1, 1), (3, 3)], [(page, number) for _, page, number in pairs])

    def test_maps_pdf_that_contains_visible_slides_only(self):
        presentation = self.make_presentation_with_hidden_slide()
        pairs = ppt2word.resolve_slide_page_pairs(presentation, 2, "sample.pptx")
        self.assertEqual([(1, 1), (2, 3)], [(page, number) for _, page, number in pairs])

    def test_get_pdf_page_count_uses_metadata(self):
        with patch.object(ppt2word, "pdfinfo_from_path", return_value={"Pages": "12"}):
            self.assertEqual(12, ppt2word.get_pdf_page_count(Path("sample.pdf")))


if __name__ == "__main__":
    unittest.main()

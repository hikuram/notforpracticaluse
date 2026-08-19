#!/usr/bin/env python3
"""Build an integrated meeting-minutes DOCX from PowerPoint files and sibling PDFs.

PowerPoint renders the source files to PDF on the Windows host. This script runs
inside the Linux container, consumes those existing PDFs, extracts text from the
original PowerPoint files, and appends the generated blocks to a DOCX template.
The source PDFs are never deleted.
"""

from __future__ import annotations

import argparse
import shutil
import sys
import tempfile
from pathlib import Path
from typing import Iterable, Iterator, Sequence

from PIL import Image
from docx import Document
from docx.document import Document as DocumentObject
from docx.enum.table import WD_ALIGN_VERTICAL, WD_ROW_HEIGHT_RULE, WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pdf2image import convert_from_path, pdfinfo_from_path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SUPPORTED_POWERPOINT_EXTENSIONS = {".pptx", ".pptm"}
DEFAULT_OUTPUT_NAME = "統合議事録ベース.docx"
DEFAULT_TEMPLATE_NAME = "header_template.docx"

# A4 portrait with 20 mm side margins leaves exactly 170 mm of body width.
PAGE_WIDTH_CM = 21.0
PAGE_HEIGHT_CM = 29.7
TOP_BOTTOM_MARGIN_CM = 1.5
LEFT_RIGHT_MARGIN_CM = 2.0
BODY_WIDTH_CM = 17.0

# A padded 16:9 image at 66 mm high is about 117.3 mm wide.  The 118.5 mm
# image cell and 0.2 mm cell margins leave only a small clearance from borders.
SLIDE_IMAGE_HEIGHT_CM = 6.6
SLIDE_IMAGE_CELL_WIDTH_CM = 11.85
SLIDE_NOTE_CELL_WIDTH_CM = BODY_WIDTH_CM - SLIDE_IMAGE_CELL_WIDTH_CM
SLIDE_IMAGE_CELL_SIDE_MARGIN_CM = 0.02
SLIDE_IMAGE_CELL_VERTICAL_MARGIN_CM = 0.02
SLIDE_ROW_MIN_HEIGHT_CM = 6.65


def sanitize_xml_text(text: str) -> str:
    """Normalize PowerPoint text so it is safe for WordprocessingML."""
    if not text:
        return ""

    # python-pptx represents PowerPoint soft line breaks as vertical tabs.
    text = text.replace("\v", "\n")

    def is_valid_xml_char(character: str) -> bool:
        codepoint = ord(character)
        return (
            codepoint in (0x09, 0x0A, 0x0D)
            or 0x20 <= codepoint <= 0xD7FF
            or 0xE000 <= codepoint <= 0xFFFD
            or 0x10000 <= codepoint <= 0x10FFFF
        )

    return "".join(character for character in text if is_valid_xml_char(character))


def find_sibling_pdf(ppt_path: Path) -> Path:
    """Return the PDF beside *ppt_path*, matching the stem case-insensitively."""
    direct_candidate = ppt_path.with_suffix(".pdf")
    if direct_candidate.is_file():
        return direct_candidate

    expected_name = f"{ppt_path.stem}.pdf".casefold()
    for candidate in ppt_path.parent.iterdir():
        if candidate.is_file() and candidate.name.casefold() == expected_name:
            return candidate

    raise FileNotFoundError(
        f"対応するPDFが見つかりません: {ppt_path.name} -> {ppt_path.stem}.pdf\n"
        "先にWindows側で pptx_to_pdf.js を実行してください。"
    )


def get_pdf_page_count(pdf_path: Path) -> int:
    """Return the PDF page count without rasterizing the complete document."""
    info = pdfinfo_from_path(str(pdf_path))
    try:
        page_count = int(info["Pages"])
    except (KeyError, TypeError, ValueError) as exc:
        raise ValueError(f"PDFページ数を取得できません: {pdf_path}") from exc
    if page_count <= 0:
        raise ValueError(f"PDFにページがありません: {pdf_path}")
    return page_count


def render_pdf_page_to_jpeg(
    pdf_path: Path,
    page_number: int,
    image_path: Path,
    *,
    dpi: int = 300,
    jpeg_quality: int = 85,
) -> None:
    """Rasterize one PDF page and save it as a JPEG."""
    images = convert_from_path(
        str(pdf_path),
        dpi=dpi,
        first_page=page_number,
        last_page=page_number,
    )
    if len(images) != 1:
        for image in images:
            image.close()
        raise ValueError(
            f"PDFページの画像化結果が1ページではありません: "
            f"{pdf_path.name} P.{page_number}"
        )

    image = images[0]
    try:
        with image.convert("RGB") as rgb_image:
            rgb_image.save(image_path, "JPEG", quality=jpeg_quality)
    finally:
        image.close()


def pad_image_to_16_9(image_path: Path, *, jpeg_quality: int = 85) -> None:
    """Contain an image in an uncropped 16:9 white canvas."""
    with Image.open(image_path) as source_image:
        with source_image.convert("RGB") as image:
            width, height = image.size
            if width <= 0 or height <= 0:
                raise ValueError(f"画像サイズが不正です: {image_path}")

            target_ratio = 16.0 / 9.0
            current_ratio = width / height

            if current_ratio < target_ratio:
                canvas_width = int(round(height * target_ratio))
                canvas_height = height
            elif current_ratio > target_ratio:
                canvas_width = width
                canvas_height = int(round(width / target_ratio))
            else:
                return

            if canvas_width == width and canvas_height == height:
                return

            with Image.new(
                "RGB",
                (canvas_width, canvas_height),
                (255, 255, 255),
            ) as canvas:
                offset_x = (canvas_width - width) // 2
                offset_y = (canvas_height - height) // 2
                canvas.paste(image, (offset_x, offset_y))
                canvas.save(image_path, "JPEG", quality=jpeg_quality)


def is_slide_visible(slide) -> bool:
    show = slide.element.get("show")
    return show not in {"0", "false"}


def normalize_extracted_text(text: str, *, line_break_separator: str) -> str:
    """Make extracted text XML-safe and compact internal line breaks."""
    text = sanitize_xml_text(text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    parts = [part.strip() for part in text.split("\n")]
    return line_break_separator.join(part for part in parts if part)


def get_paragraph_text(paragraph, *, line_break_separator: str) -> str:
    """Return compact paragraph text with run hyperlink targets appended."""
    run_by_element = {run._r: run for run in paragraph.runs}
    chunks: list[str] = []
    active_hyperlink: str | None = None

    def flush_hyperlink() -> None:
        nonlocal active_hyperlink
        if active_hyperlink:
            chunks.append(f"（{active_hyperlink}）")
            active_hyperlink = None

    for element in paragraph._p.content_children:
        local_name = element.tag.rsplit("}", 1)[-1]

        if local_name == "r":
            run = run_by_element.get(element)
            hyperlink = ""
            if run is not None:
                hyperlink = sanitize_xml_text(run.hyperlink.address or "").strip()

            if active_hyperlink and hyperlink != active_hyperlink:
                flush_hyperlink()

            chunks.append(element.text or "")
            if hyperlink:
                active_hyperlink = hyperlink
            continue

        flush_hyperlink()
        if local_name == "br":
            chunks.append("\n")
        else:
            chunks.append(element.text or "")

    flush_hyperlink()
    return normalize_extracted_text(
        "".join(chunks),
        line_break_separator=line_break_separator,
    )


def get_text_frame_text(
    text_frame,
    *,
    paragraph_separator: str = "\t",
    line_break_separator: str = "\t",
) -> str:
    """Return compact text from one text frame while preserving run joins."""
    paragraphs: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = get_paragraph_text(
            paragraph,
            line_break_separator=line_break_separator,
        )
        if text:
            paragraphs.append(text)
    return paragraph_separator.join(paragraphs)


def iter_table_text(table) -> Iterator[str]:
    """Yield one compact searchable line per PowerPoint table row."""
    seen_cells: set[object] = set()
    for row in table.rows:
        cell_texts: list[str] = []
        for cell in row.cells:
            cell_key = cell._tc
            if cell_key in seen_cells:
                continue
            seen_cells.add(cell_key)
            cell_texts.append(
                get_text_frame_text(
                    cell.text_frame,
                    paragraph_separator=" ",
                    line_break_separator=" ",
                )
            )
        if any(cell_texts):
            yield "\t".join(cell_texts).strip("\t")


def iter_shape_text(shape) -> Iterator[str]:
    """Yield searchable text from a shape, including groups and tables."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            yield from iter_shape_text(child_shape)
        return

    if getattr(shape, "has_table", False):
        yield from iter_table_text(shape.table)
        return

    if getattr(shape, "has_text_frame", False):
        text = get_text_frame_text(shape.text_frame)
        if text:
            yield text


def extract_text_from_slide(slide) -> str:
    """Extract compact searchable slide text in shape order."""
    shape_texts: list[str] = []
    for shape in slide.shapes:
        shape_texts.extend(iter_shape_text(shape))
    return "\n".join(shape_texts)


def get_or_add_child(parent, tag: str) -> OxmlElement:
    child = parent.find(qn(tag))
    if child is None:
        child = OxmlElement(tag)
        parent.append(child)
    return child


def set_contextual_spacing(element) -> None:
    """Enable Word's 'do not add space between same-style paragraphs' flag."""
    paragraph_properties = element.get_or_add_pPr()
    if paragraph_properties.find(qn("w:contextualSpacing")) is None:
        paragraph_properties.append(OxmlElement("w:contextualSpacing"))


def disable_snap_to_grid(element) -> None:
    """Disable Word's 'snap to document grid' paragraph option."""
    paragraph_properties = element.get_or_add_pPr()
    snap_to_grid = paragraph_properties.find(qn("w:snapToGrid"))
    if snap_to_grid is None:
        snap_to_grid = OxmlElement("w:snapToGrid")
        paragraph_properties.append(snap_to_grid)
    snap_to_grid.set(qn("w:val"), "0")


def enable_track_revisions(document: DocumentObject) -> None:
    """Start the generated document with Track Changes enabled for all editors."""
    settings = document.settings._element
    track_revisions = settings.find(qn("w:trackRevisions"))
    if track_revisions is None:
        track_revisions = OxmlElement("w:trackRevisions")
        settings.append(track_revisions)
    track_revisions.set(qn("w:val"), "true")


def set_outline_level(paragraph, level: int) -> None:
    paragraph_properties = paragraph._p.get_or_add_pPr()
    outline = paragraph_properties.find(qn("w:outlineLvl"))
    if outline is None:
        outline = OxmlElement("w:outlineLvl")
        paragraph_properties.append(outline)
    outline.set(qn("w:val"), str(level))


def apply_collapse_hack(paragraph) -> None:
    try:
        paragraph_properties = paragraph._p.get_or_add_pPr()
        existing = paragraph_properties.find(
            "{http://schemas.microsoft.com/office/word/2012/wordml}collapsed"
        )
        if existing is not None:
            paragraph_properties.remove(existing)
        collapsed_xml = (
            '<w15:collapsed w:val="1" '
            'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
            'xmlns:w15="http://schemas.microsoft.com/office/word/2012/wordml"/>'
        )
        paragraph_properties.append(parse_xml(collapsed_xml))
    except Exception:
        # The document remains usable even if a Word version ignores the hack.
        pass


def iter_table_paragraphs(table) -> Iterator:
    seen_cells: set[int] = set()
    for row in table.rows:
        for cell in row.cells:
            cell_key = id(cell._tc)
            if cell_key in seen_cells:
                continue
            seen_cells.add(cell_key)
            yield from cell.paragraphs
            for nested_table in cell.tables:
                yield from iter_table_paragraphs(nested_table)


def iter_all_paragraphs(document: DocumentObject) -> Iterator:
    yield from document.paragraphs
    for table in document.tables:
        yield from iter_table_paragraphs(table)


def compact_paragraph(paragraph, *, before_pt: float | None = None) -> None:
    if before_pt is not None:
        paragraph.paragraph_format.space_before = Pt(before_pt)
    paragraph.paragraph_format.space_after = Pt(0)
    set_contextual_spacing(paragraph._p)
    disable_snap_to_grid(paragraph._p)


def set_style_font_to_meiryo(document: DocumentObject) -> None:
    target_styles = [
        "Normal",
        "Title",
        "Heading 1",
        "Heading 2",
        "Heading 3",
        "Table Grid",
    ]
    for style_name in target_styles:
        try:
            style = document.styles[style_name]
        except KeyError:
            continue

        style.font.name = "メイリオ"
        if style.type == WD_STYLE_TYPE.PARAGRAPH:
            disable_snap_to_grid(style.element)
        run_properties = style.element.get_or_add_rPr()
        font_elements = run_properties.xpath("w:rFonts")
        if font_elements:
            font_elements[0].set(qn("w:eastAsia"), "メイリオ")
        else:
            fonts = OxmlElement("w:rFonts")
            fonts.set(qn("w:eastAsia"), "メイリオ")
            run_properties.append(fonts)


def ensure_heading_styles(document: DocumentObject) -> None:
    """Create Heading 1-3 when the supplied form template omits them."""
    normal_style = document.styles["Normal"]
    sizes = {"Heading 1": 11.0, "Heading 2": 9.5, "Heading 3": 8.5}
    for index, style_name in enumerate(("Heading 1", "Heading 2", "Heading 3"), start=1):
        try:
            style = document.styles[style_name]
        except KeyError:
            style = document.styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
            style.base_style = normal_style
        style.font.name = "メイリオ"
        style.font.size = Pt(sizes[style_name])
        style.font.bold = True
        style.hidden = False
        style.unhide_when_used = False
        style.quick_style = True
        style.priority = 8 + index


def configure_heading_style(document: DocumentObject, style_name: str, level: int) -> None:
    try:
        style = document.styles[style_name]
    except KeyError:
        return

    style.paragraph_format.space_before = Pt(0)
    style.paragraph_format.space_after = Pt(0)
    style.paragraph_format.keep_with_next = True
    style.paragraph_format.keep_together = True

    paragraph_properties = style.element.get_or_add_pPr()
    outline = paragraph_properties.find(qn("w:outlineLvl"))
    if outline is None:
        outline = OxmlElement("w:outlineLvl")
        paragraph_properties.append(outline)
    outline.set(qn("w:val"), str(level))

    if paragraph_properties.find(qn("w:contextualSpacing")) is None:
        paragraph_properties.append(OxmlElement("w:contextualSpacing"))
    disable_snap_to_grid(style.element)


def configure_template_document(document: DocumentObject) -> None:
    """Keep the supplied template as the document base and normalize essentials."""
    for section in document.sections:
        section.page_width = Cm(PAGE_WIDTH_CM)
        section.page_height = Cm(PAGE_HEIGHT_CM)
        section.top_margin = Cm(TOP_BOTTOM_MARGIN_CM)
        section.bottom_margin = Cm(TOP_BOTTOM_MARGIN_CM)
        section.left_margin = Cm(LEFT_RIGHT_MARGIN_CM)
        section.right_margin = Cm(LEFT_RIGHT_MARGIN_CM)

    ensure_heading_styles(document)
    set_style_font_to_meiryo(document)

    try:
        normal_style = document.styles["Normal"]
        normal_style.paragraph_format.space_after = Pt(0)
        normal_ppr = normal_style.element.get_or_add_pPr()
        if normal_ppr.find(qn("w:contextualSpacing")) is None:
            normal_ppr.append(OxmlElement("w:contextualSpacing"))
        disable_snap_to_grid(normal_style.element)
    except KeyError:
        pass

    configure_heading_style(document, "Heading 1", 0)
    configure_heading_style(document, "Heading 2", 1)
    configure_heading_style(document, "Heading 3", 2)

    # Retain the template's direct formatting while ensuring that Word's
    # paragraph-after-spacing removal remains active throughout the document.
    for paragraph in iter_all_paragraphs(document):
        compact_paragraph(paragraph)

    # The attached template is already aligned to the 20 mm page margins.  Make
    # that explicit so copied/edited versions cannot retain a small negative
    # table indent from an older template.
    for table in document.tables:
        set_table_left_alignment(table)


def set_table_left_alignment(table) -> None:
    table.alignment = WD_TABLE_ALIGNMENT.LEFT
    table_properties = table._tbl.tblPr

    table_indent = table_properties.find(qn("w:tblInd"))
    if table_indent is None:
        table_indent = OxmlElement("w:tblInd")
        table_properties.append(table_indent)
    table_indent.set(qn("w:w"), "0")
    table_indent.set(qn("w:type"), "dxa")

    justification = table_properties.find(qn("w:jc"))
    if justification is None:
        justification = OxmlElement("w:jc")
        table_properties.append(justification)
    justification.set(qn("w:val"), "left")


def set_table_border_color(table, color: str = "D3D3D3") -> None:
    table_properties = table._tbl.tblPr
    existing_borders = table_properties.find(qn("w:tblBorders"))
    if existing_borders is not None:
        table_properties.remove(existing_borders)

    borders = OxmlElement("w:tblBorders")
    for border_name in ("top", "left", "bottom", "right", "insideH", "insideV"):
        border = OxmlElement(f"w:{border_name}")
        border.set(qn("w:val"), "single")
        border.set(qn("w:sz"), "4")
        border.set(qn("w:space"), "0")
        border.set(qn("w:color"), color)
        borders.append(border)
    table_properties.append(borders)


def set_cell_margins(
    cell,
    *,
    top_cm: float,
    start_cm: float,
    bottom_cm: float,
    end_cm: float,
) -> None:
    cell_properties = cell._tc.get_or_add_tcPr()
    margins = cell_properties.find(qn("w:tcMar"))
    if margins is None:
        margins = OxmlElement("w:tcMar")
        cell_properties.append(margins)

    for side, value_cm in (
        ("top", top_cm),
        ("start", start_cm),
        ("bottom", bottom_cm),
        ("end", end_cm),
    ):
        element = margins.find(qn(f"w:{side}"))
        if element is None:
            element = OxmlElement(f"w:{side}")
            margins.append(element)
        element.set(qn("w:w"), str(round(value_cm * 567)))
        element.set(qn("w:type"), "dxa")


def prevent_row_split(row) -> None:
    row_properties = row._tr.get_or_add_trPr()
    if row_properties.find(qn("w:cantSplit")) is None:
        row_properties.append(OxmlElement("w:cantSplit"))


def enforce_absolute_table_width(table, column_widths_cm: Sequence[float]) -> None:
    if len(column_widths_cm) != len(table.columns):
        raise ValueError("列幅指定と表の列数が一致しません。")

    table.autofit = False
    set_table_left_alignment(table)
    table_properties = table._tbl.tblPr

    layout = table_properties.find(qn("w:tblLayout"))
    if layout is None:
        layout = OxmlElement("w:tblLayout")
        table_properties.append(layout)
    layout.set(qn("w:type"), "fixed")

    table_width = table_properties.find(qn("w:tblW"))
    if table_width is None:
        table_width = OxmlElement("w:tblW")
        table_properties.append(table_width)
    table_width.set(qn("w:w"), str(round(sum(column_widths_cm) * 567)))
    table_width.set(qn("w:type"), "dxa")

    grid = table._tbl.tblGrid
    for grid_column, width_cm in zip(grid.gridCol_lst, column_widths_cm):
        grid_column.set(qn("w:w"), str(round(width_cm * 567)))

    for column, width_cm in zip(table.columns, column_widths_cm):
        column.width = Cm(width_cm)

    for row in table.rows:
        for cell, width_cm in zip(row.cells, column_widths_cm):
            cell.width = Cm(width_cm)
            cell_properties = cell._tc.get_or_add_tcPr()
            cell_width = cell_properties.find(qn("w:tcW"))
            if cell_width is None:
                cell_width = OxmlElement("w:tcW")
                cell_properties.append(cell_width)
            cell_width.set(qn("w:w"), str(round(width_cm * 567)))
            cell_width.set(qn("w:type"), "dxa")


def resolve_slide_page_pairs(
    presentation: Presentation,
    pdf_page_count: int,
    source_name: str,
) -> list[tuple[object, int, int]]:
    """Map visible slides to 1-based PDF pages for either export mode."""
    all_slides = list(presentation.slides)
    visible_entries = [
        (slide, original_index)
        for original_index, slide in enumerate(all_slides, start=1)
        if is_slide_visible(slide)
    ]

    if pdf_page_count == len(all_slides):
        return [
            (slide, original_index, original_index)
            for slide, original_index in visible_entries
        ]

    if pdf_page_count == len(visible_entries):
        return [
            (slide, pdf_page_number, original_index)
            for pdf_page_number, (slide, original_index)
            in enumerate(visible_entries, start=1)
        ]

    raise ValueError(
        f"{source_name}: PDFページ数 ({pdf_page_count}) が、"
        f"全スライド数 ({len(all_slides)}) または表示スライド数 "
        f"({len(visible_entries)}) と一致しません。"
    )


def add_file_header(document: DocumentObject, source_name: str) -> None:
    file_heading = document.add_paragraph(f"▼ 資料: {source_name}")
    file_heading.style = "Heading 1"
    compact_paragraph(file_heading)
    file_heading.paragraph_format.keep_with_next = True
    set_outline_level(file_heading, 0)
    for run in file_heading.runs:
        run.font.color.rgb = RGBColor(0, 51, 153)
        run.font.size = Pt(11)
        run.bold = True

    file_table = document.add_table(rows=1, cols=2)
    file_table.style = "Table Grid"
    enforce_absolute_table_width(file_table, (8.5, 8.5))

    row = file_table.rows[0]
    row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
    row.height = Cm(6.0)
    prevent_row_split(row)

    file_table.cell(0, 0).text = "特記事項："
    for cell in row.cells:
        cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP
        set_cell_margins(cell, top_cm=0.12, start_cm=0.15, bottom_cm=0.12, end_cm=0.15)
        for paragraph in cell.paragraphs:
            compact_paragraph(paragraph)

    spacer = document.add_paragraph()
    compact_paragraph(spacer)


def add_slide_block(
    document: DocumentObject,
    *,
    source_name: str,
    slide,
    image_path: Path,
    original_slide_number: int,
) -> None:
    slide_text = extract_text_from_slide(slide)

    slide_heading = document.add_paragraph(
        f"■ {source_name} - P.{original_slide_number}"
    )
    slide_heading.style = "Heading 2"
    compact_paragraph(slide_heading)
    slide_heading.paragraph_format.keep_with_next = True
    set_outline_level(slide_heading, 1)
    for run in slide_heading.runs:
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor(0, 51, 153)
        run.bold = True

    slide_table = document.add_table(rows=1, cols=2)
    slide_table.style = "Table Grid"
    enforce_absolute_table_width(
        slide_table,
        (SLIDE_IMAGE_CELL_WIDTH_CM, SLIDE_NOTE_CELL_WIDTH_CM),
    )
    set_table_border_color(slide_table, "B0B0B0")

    row = slide_table.rows[0]
    row.height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
    row.height = Cm(SLIDE_ROW_MIN_HEIGHT_CM)
    prevent_row_split(row)

    image_cell = row.cells[0]
    note_cell = row.cells[1]
    image_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
    note_cell.vertical_alignment = WD_ALIGN_VERTICAL.TOP

    set_cell_margins(
        image_cell,
        top_cm=SLIDE_IMAGE_CELL_VERTICAL_MARGIN_CM,
        start_cm=SLIDE_IMAGE_CELL_SIDE_MARGIN_CM,
        bottom_cm=SLIDE_IMAGE_CELL_VERTICAL_MARGIN_CM,
        end_cm=SLIDE_IMAGE_CELL_SIDE_MARGIN_CM,
    )
    set_cell_margins(
        note_cell,
        top_cm=0.12,
        start_cm=0.15,
        bottom_cm=0.12,
        end_cm=0.15,
    )

    image_paragraph = image_cell.paragraphs[0]
    compact_paragraph(image_paragraph)
    image_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    image_paragraph.paragraph_format.left_indent = Cm(0)
    image_paragraph.paragraph_format.right_indent = Cm(0)
    image_paragraph.add_run().add_picture(
        str(image_path),
        height=Cm(SLIDE_IMAGE_HEIGHT_CM),
    )

    note_paragraph = note_cell.paragraphs[0]
    compact_paragraph(note_paragraph)

    if slide_text:
        summary_heading = document.add_paragraph("▶ 抽出テキスト（検索・コピペ用）")
        summary_heading.style = "Heading 3"
        compact_paragraph(summary_heading)
        summary_heading.paragraph_format.keep_with_next = True
        set_outline_level(summary_heading, 2)
        apply_collapse_hack(summary_heading)
        for run in summary_heading.runs:
            run.font.size = Pt(8.5)
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.bold = True

        text_paragraph = document.add_paragraph(slide_text)
        text_paragraph.style = "Normal"
        compact_paragraph(text_paragraph)
        for run in text_paragraph.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(80, 80, 80)

    spacer = document.add_paragraph()
    compact_paragraph(spacer)


def build_minutes_base(
    ppt_paths: Sequence[Path],
    output_word_path: Path,
    *,
    template_path: Path,
    dpi: int = 300,
    jpeg_quality: int = 85,
    keep_images: bool = False,
) -> None:
    if output_word_path.resolve() == template_path.resolve():
        raise ValueError("出力先にテンプレート自身は指定できません。")

    document = Document(str(template_path))
    configure_template_document(document)
    enable_track_revisions(document)

    temp_parent = output_word_path.parent
    temp_parent.mkdir(parents=True, exist_ok=True)
    temp_path = Path(tempfile.mkdtemp(prefix="_ppt_temp_v5_", dir=temp_parent))

    try:
        for file_index, ppt_path in enumerate(ppt_paths, start=1):
            source_name = ppt_path.name
            pdf_path = find_sibling_pdf(ppt_path)
            item_temp_dir = temp_path / f"{file_index:03d}_{ppt_path.stem}"
            item_temp_dir.mkdir(parents=True, exist_ok=True)

            print(f"\n処理中: {source_name}")
            print(f"  PDF: {pdf_path.name}")

            presentation = Presentation(str(ppt_path))
            pdf_page_count = get_pdf_page_count(pdf_path)
            slide_page_pairs = resolve_slide_page_pairs(
                presentation,
                pdf_page_count,
                source_name,
            )

            add_file_header(document, source_name)
            for slide, pdf_page_number, original_slide_number in slide_page_pairs:
                image_path = item_temp_dir / f"slide_{original_slide_number:04d}.jpg"
                render_pdf_page_to_jpeg(
                    pdf_path,
                    pdf_page_number,
                    image_path,
                    dpi=dpi,
                    jpeg_quality=jpeg_quality,
                )
                pad_image_to_16_9(image_path, jpeg_quality=jpeg_quality)
                add_slide_block(
                    document,
                    source_name=source_name,
                    slide=slide,
                    image_path=image_path,
                    original_slide_number=original_slide_number,
                )
                if not keep_images:
                    image_path.unlink(missing_ok=True)

        # Reapply the compact paragraph setting after all generated content has
        # been added, without changing the template's other direct formatting.
        for paragraph in iter_all_paragraphs(document):
            paragraph.paragraph_format.space_after = Pt(0)
            set_contextual_spacing(paragraph._p)
            disable_snap_to_grid(paragraph._p)

        print(f"\n保存中: {output_word_path.name}")
        document.save(str(output_word_path))
    finally:
        if keep_images:
            print(f"中間画像を保持しました: {temp_path}")
        else:
            shutil.rmtree(temp_path, ignore_errors=True)

    print("完了: 統合議事録ベースを生成しました。")


def discover_powerpoint_files(directory: Path) -> list[Path]:
    return sorted(
        (
            path.resolve()
            for path in directory.iterdir()
            if path.is_file() and path.suffix.casefold() in SUPPORTED_POWERPOINT_EXTENSIONS
        ),
        key=lambda path: path.name.casefold(),
    )


def validate_inputs(raw_paths: Iterable[str], current_directory: Path) -> list[Path]:
    raw_paths = list(raw_paths)
    if not raw_paths:
        discovered = discover_powerpoint_files(current_directory)
        if not discovered:
            raise FileNotFoundError(
                f"カレントディレクトリにPPTX/PPTMがありません: {current_directory}"
            )
        return discovered

    validated: list[Path] = []
    seen: set[Path] = set()
    for raw_path in raw_paths:
        path = Path(raw_path).expanduser().resolve()
        if not path.is_file():
            raise FileNotFoundError(f"PowerPointファイルが見つかりません: {raw_path}")
        if path.suffix.casefold() not in SUPPORTED_POWERPOINT_EXTENSIONS:
            raise ValueError(f"未対応のPowerPoint形式です: {path.name}")
        if path not in seen:
            seen.add(path)
            validated.append(path)
    return validated


def resolve_template_path(raw_template_path: str | None, current_directory: Path) -> Path:
    if raw_template_path:
        template_path = Path(raw_template_path).expanduser()
        if not template_path.is_absolute():
            template_path = current_directory / template_path
    else:
        template_path = Path(__file__).resolve().with_name(DEFAULT_TEMPLATE_NAME)

    template_path = template_path.resolve()
    if not template_path.is_file():
        raise FileNotFoundError(f"DOCXテンプレートが見つかりません: {template_path}")
    return template_path


def parse_arguments(argv: Sequence[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Windows PowerPointで生成済みの同名PDFを使い、"
            "添付様式をベースに統合議事録DOCXを生成します。"
        )
    )
    parser.add_argument(
        "ppt_files",
        nargs="*",
        help="処理するPPTX/PPTM。省略時はカレントディレクトリを走査します。",
    )
    parser.add_argument(
        "-o",
        "--output",
        default=DEFAULT_OUTPUT_NAME,
        help=f"出力DOCXパス（既定: {DEFAULT_OUTPUT_NAME}）",
    )
    parser.add_argument(
        "--template",
        "--header-template",
        "--header-block",
        dest="template",
        default=None,
        help=(
            "文書全体のベースにするDOCX。省略時はイメージ内蔵の "
            f"{DEFAULT_TEMPLATE_NAME} を使用します。"
        ),
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=300,
        help="PDF画像化の解像度（既定: 300）",
    )
    parser.add_argument(
        "--jpeg-quality",
        type=int,
        default=85,
        help="中間JPEG品質 1-95（既定: 85）",
    )
    parser.add_argument(
        "--keep-images",
        action="store_true",
        help="デバッグ用に中間JPEGを残します。PDFは常に残ります。",
    )
    return parser.parse_args(argv)


def main(argv: Sequence[str] | None = None) -> int:
    args = parse_arguments(argv if argv is not None else sys.argv[1:])

    if args.dpi <= 0:
        print("エラー: --dpi は正の整数で指定してください。", file=sys.stderr)
        return 2
    if not 1 <= args.jpeg_quality <= 95:
        print("エラー: --jpeg-quality は1から95で指定してください。", file=sys.stderr)
        return 2

    current_directory = Path.cwd().resolve()
    try:
        ppt_paths = validate_inputs(args.ppt_files, current_directory)
        template_path = resolve_template_path(args.template, current_directory)

        output_path = Path(args.output).expanduser()
        if not output_path.is_absolute():
            output_path = current_directory / output_path
        output_path = output_path.resolve()

        for ppt_path in ppt_paths:
            find_sibling_pdf(ppt_path)

        build_minutes_base(
            ppt_paths,
            output_path,
            template_path=template_path,
            dpi=args.dpi,
            jpeg_quality=args.jpeg_quality,
            keep_images=args.keep_images,
        )
    except (FileNotFoundError, ValueError, OSError) as error:
        print(f"エラー: {error}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
